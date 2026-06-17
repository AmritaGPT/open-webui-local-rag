import os
import sys
import pandas as pd

# Add the tools directory to path so we can import the scripts
sys.path.append(os.path.join(os.path.dirname(__file__), 'open_webui_tools'))

# Import our tools
from local_document_parser import Tools as ParserTools
from spreadsheet_query import Tools as SpreadsheetTools

def generate_sample_data():
    os.makedirs('sample_data', exist_ok=True)
    print("Generating sample files in './sample_data/'...")

    # 1. Generate Excel (sales.xlsx)
    sales_data = {
        'Product': ['Laptop', 'Mouse', 'Keyboard', 'Monitor', 'USB Cable'],
        'Category': ['Electronics', 'Accessories', 'Accessories', 'Electronics', 'Accessories'],
        'Units Sold': [120, 450, 300, 80, 1500],
        'UnitPrice': [999.99, 25.50, 45.00, 249.99, 8.99],
        'Region': ['North', 'East', 'West', 'South', 'East']
    }
    df = pd.DataFrame(sales_data)
    df['Revenue'] = df['Units Sold'] * df['UnitPrice']
    df.to_excel('sample_data/sales.xlsx', index=False)
    print("-> Generated sample_data/sales.xlsx")

    # 1b. Generate Multi-sheet Excel (multisheet.xlsx)
    q1_data = {'Product': ['Laptop', 'Mouse'], 'Revenue': [100000, 12000]}
    q2_data = {'Product': ['Laptop', 'Mouse'], 'Revenue': [130000, 15000]}
    
    with pd.ExcelWriter('sample_data/multisheet.xlsx') as writer:
        pd.DataFrame(q1_data).to_excel(writer, sheet_name='Q1_Sales', index=False)
        pd.DataFrame(q2_data).to_excel(writer, sheet_name='Q2_Sales', index=False)
    print("-> Generated sample_data/multisheet.xlsx")

    # 2. Generate Word (contract.docx)
    try:
        import docx
        doc = docx.Document()
        doc.add_heading('Service Agreement', 0)
        
        doc.add_paragraph('This Service Agreement ("Agreement") is made effective as of June 17, 2026.')
        
        doc.add_heading('1. Scope of Work', level=1)
        doc.add_paragraph('The contractor agrees to provide technical support and development services as outlined in the table below.')
        
        # Add a table
        table = doc.add_table(rows=3, cols=3)
        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = 'Task ID'
        hdr_cells[1].text = 'Description'
        hdr_cells[2].text = 'Estimated Hours'
        
        row1_cells = table.rows[1].cells
        row1_cells[0].text = 'T001'
        row1_cells[1].text = 'API Integration'
        row1_cells[2].text = '40'
        
        row2_cells = table.rows[2].cells
        row2_cells[0].text = 'T002'
        row2_cells[1].text = 'UI Testing'
        row2_cells[2].text = '20'
        
        doc.save('sample_data/contract.docx')
        print("-> Generated sample_data/contract.docx")
    except Exception as e:
        print(f"Failed to generate Word document: {e}")

    # 3. Generate PowerPoint (slides.pptx)
    try:
        from pptx import Presentation
        prs = Presentation()
        
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = "Q2 Performance Review"
        subtitle.text = "Presented by Advanced Agentic Team"
        
        # Slide 2: Table slide
        slide_layout2 = prs.slide_layouts[5] # title only
        slide2 = prs.slides.add_slide(slide_layout2)
        slide2.shapes.title.text = "Regional Performance Matrix"
        
        # Add table
        rows, cols = 3, 3
        left = top = width = height = docx.shared.Inches(2.0)
        table_shape = slide2.shapes.add_table(rows, cols, left, top, width, height)
        table = table_shape.table
        
        # Set headers
        table.cell(0, 0).text = 'Region'
        table.cell(0, 1).text = 'Q1 Revenue ($)'
        table.cell(0, 2).text = 'Q2 Revenue ($)'
        
        # Data
        table.cell(1, 0).text = 'North'
        table.cell(1, 1).text = '500,000'
        table.cell(1, 2).text = '620,000'
        
        table.cell(2, 0).text = 'South'
        table.cell(2, 1).text = '350,000'
        table.cell(2, 2).text = '410,000'
        
        prs.save('sample_data/slides.pptx')
        print("-> Generated sample_data/slides.pptx")
    except Exception as e:
        print(f"Failed to generate PowerPoint: {e}")

    # 4. Generate PDF (report.pdf)
    try:
        from reportlab.lib.pagesizes import letter
        from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.lib import colors
        
        doc = SimpleDocTemplate("sample_data/report.pdf", pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        story.append(Paragraph("Annual Financial Report", styles['Title']))
        story.append(Spacer(1, 20))
        story.append(Paragraph("This report highlights the key metrics for the financial year ending June 2026. The table below represents the performance of different business units.", styles['Normal']))
        story.append(Spacer(1, 10))
        
        data = [
            ['Unit', 'Direct Cost ($)', 'Net Earnings ($)'],
            ['Cloud Services', '1,200,000', '4,500,000'],
            ['Enterprise Software', '800,000', '3,100,000'],
            ['Hardware Support', '600,000', '1,200,000']
        ]
        
        t = Table(data)
        t.setStyle(TableStyle([
            ('BACKGROUND', (0,0), (-1,0), colors.grey),
            ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
            ('ALIGN', (0,0), (-1,-1), 'CENTER'),
            ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
            ('BOTTOMPADDING', (0,0), (-1,0), 8),
            ('BACKGROUND', (0,1), (-1,-1), colors.beige),
            ('GRID', (0,0), (-1,-1), 1, colors.black),
        ]))
        story.append(t)
        doc.build(story)
        print("-> Generated sample_data/report.pdf")
    except Exception as e:
        print(f"Failed to generate PDF: {e}")

def run_tests():
    print("\n" + "="*50)
    print("RUNNING PARSING TESTS ON PROGRAMMATICALLY GENERATED DATA")
    print("="*50)
    
    parser = ParserTools()
    spreadsheet_tool = SpreadsheetTools()
    
    # Test 1: PDF Parse
    print("\n--- Test 1: PDF Local Parser ---")
    pdf_out = parser.parse_local_document("sample_data/report.pdf")
    print(pdf_out[:1000]) # Print first 1000 characters
    
    # Test 2: Word Parse
    print("\n--- Test 2: Word Local Parser ---")
    docx_out = parser.parse_local_document("sample_data/contract.docx")
    print(docx_out)
    
    # Test 3: PowerPoint Parse
    print("\n--- Test 3: PowerPoint Local Parser ---")
    pptx_out = parser.parse_local_document("sample_data/slides.pptx")
    print(pptx_out)

    # Test 4: Excel Pandas Code Execution
    print("\n--- Test 4: Excel Spreadsheet Analyzer (Local Pandas) ---")
    # Let's run a pandas code snippet to find total revenue
    code = "result = df['Revenue'].sum()"
    excel_out = spreadsheet_tool.run_pandas_code("sample_data/sales.xlsx", code)
    print(f"Total Revenue from Excel calculation: ${float(excel_out):,.2f}")

    # Let's run a grouping query
    code_group = "result = df.groupby('Product')['Revenue'].sum()"
    excel_out_group = spreadsheet_tool.run_pandas_code("sample_data/sales.xlsx", code_group)
    print(f"Revenue by Product:\n{excel_out_group}")

    # Test 5: Multi-sheet Excel Edge Case
    print("\n--- Test 5: Multi-sheet Excel Edge Case ---")
    # Retrieve schema first
    schema = spreadsheet_tool.get_spreadsheet_schema("sample_data/multisheet.xlsx")
    print(schema)
    
    # Run code to sum revenues from both sheets
    code_multisheet = "result = sheets['Q1_Sales']['Revenue'].sum() + sheets['Q2_Sales']['Revenue'].sum()"
    excel_out_multisheet = spreadsheet_tool.run_pandas_code("sample_data/multisheet.xlsx", code_multisheet)
    print(f"Combined Q1 + Q2 Revenue: ${float(excel_out_multisheet):,.2f}")

if __name__ == '__main__':
    generate_sample_data()
    run_tests()
