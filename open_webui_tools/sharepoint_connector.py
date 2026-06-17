"""
title: Local SharePoint Connector & Parser
author: Advanced Agentic RAG Team
version: 1.0.0
description: Integrates with SharePoint (either via a local OneDrive/SharePoint synced directory or directly via Microsoft Graph Cloud API) to search, download, and parse files locally and offline.
requirements: msal, requests, pdfplumber, python-docx, python-pptx, pytesseract, Pillow
"""

import os
import glob
from typing import Dict, List, Any
from pydantic import BaseModel, Field
import requests

class Tools:
    class Valves(BaseModel):
        SHAREPOINT_LOCAL_PATH: str = Field(
            default="",
            description="Absolute path to your synced SharePoint or OneDrive folder on this computer. If set, searches will run locally (Recommended for offline setup)."
        )
        TENANT_ID: str = Field(
            default="",
            description="Azure AD Tenant ID (Required only for Cloud API connection)."
        )
        CLIENT_ID: str = Field(
            default="",
            description="Azure AD Client ID (Required only for Cloud API connection)."
        )
        CLIENT_SECRET: str = Field(
            default="",
            description="Azure AD Client Secret (Required only for Cloud API connection)."
        )
        SITE_ID: str = Field(
            default="",
            description="SharePoint Site ID or Site Name (Required only for Cloud API connection)."
        )
        AUTHORIZED_EMAILS: str = Field(
            default="",
            description="Comma-separated list of user emails authorized to use this tool (e.g. 'alice@company.com, bob@company.com'). Leave blank to allow all users."
        )

    def __init__(self):
        self.valves = self.Valves()

    def _is_authorized(self, __user__: dict = None) -> bool:
        if not self.valves.AUTHORIZED_EMAILS:
            return True
        if not __user__:
            return False
        # Admins always bypass checks
        if __user__.get("role") == "admin":
            return True
        user_email = __user__.get("email", "").lower().strip()
        auth_list = [e.lower().strip() for e in self.valves.AUTHORIZED_EMAILS.split(",") if e.strip()]
        return user_email in auth_list

    def search_documents(self, query: str, __user__: dict = None) -> str:
        """
        Searches the SharePoint repository (synced folder or cloud) for documents matching the query keywords.
        Returns a list of matching filenames, paths, and document sizes.
        
        :param query: The search term or keywords (e.g., 'NDA template', 'Q3 financial report').
        :return: A list of matching files formatted as markdown.
        """
        if not self._is_authorized(__user__):
            return "Error: You are not authorized to use the SharePoint document search tool."

        # Scenario A: Local Synced Folder (Primary offline method)
        if self.valves.SHAREPOINT_LOCAL_PATH:
            if not os.path.exists(self.valves.SHAREPOINT_LOCAL_PATH):
                return f"Error: Local SharePoint folder path '{self.valves.SHAREPOINT_LOCAL_PATH}' does not exist."
            
            # Search recursively using globbing
            search_pattern = os.path.join(self.valves.SHAREPOINT_LOCAL_PATH, "**", f"*{query}*")
            matches = glob.glob(search_pattern, recursive=True)
            
            if not matches:
                # If no direct match, try searching files containing query in their names
                search_pattern = os.path.join(self.valves.SHAREPOINT_LOCAL_PATH, "**", "*")
                all_files = glob.glob(search_pattern, recursive=True)
                matches = [f for f in all_files if os.path.isfile(f) and query.lower() in os.path.basename(f).lower()]
                
            if not matches:
                return f"No documents matching '{query}' were found in local path '{self.valves.SHAREPOINT_LOCAL_PATH}'."
                
            results = ["### Synced SharePoint Files Found:\n"]
            for filepath in matches[:10]: # Limit to top 10 matches
                rel_path = os.path.relpath(filepath, self.valves.SHAREPOINT_LOCAL_PATH)
                size_kb = os.path.getsize(filepath) / 1024
                results.append(f"- **Filename:** `{os.path.basename(filepath)}`\n  - **Path:** `{rel_path}`\n  - **Size:** {size_kb:.1f} KB")
            return "\n".join(results)

        # Scenario B: Cloud SharePoint Connection (MS Graph API)
        elif self.valves.CLIENT_ID and self.valves.CLIENT_SECRET and self.valves.TENANT_ID:
            try:
                import msal
                token = self._get_microsoft_token()
                headers = {"Authorization": f"Bearer {token}"}
                
                # If site ID is not specified, search user's OneDrive/SharePoint root
                if self.valves.SITE_ID:
                    url = f"https://graph.microsoft.com/v1.0/sites/{self.valves.SITE_ID}/drive/root/search(q='{query}')"
                else:
                    url = f"https://graph.microsoft.com/v1.0/me/drive/root/search(q='{query}')"
                    
                response = requests.get(url, headers=headers)
                response.raise_for_status()
                items = response.json().get('value', [])
                
                if not items:
                    return f"No documents matching '{query}' were found in Cloud SharePoint."
                    
                results = ["### Cloud SharePoint Files Found:\n"]
                for item in items[:10]:
                    results.append(f"- **Filename:** `{item.get('name')}`\n  - **ID:** `{item.get('id')}`\n  - **Web URL:** {item.get('webUrl')}\n  - **Size:** {item.get('size', 0)/1024:.1f} KB")
                return "\n".join(results)
            except Exception as e:
                return f"Cloud SharePoint search failed: {str(e)}"
                
        else:
            return "Error: SharePoint connection not configured. Please specify 'SHAREPOINT_LOCAL_PATH' for offline files or Microsoft Graph credentials in the tool settings (Valves)."

    def read_and_parse_document(self, identifier: str, __user__: dict = None) -> str:
        """
        Downloads (if cloud) and parses a SharePoint document using layout-aware local parsers.
        Accepts either a relative local path or a cloud File ID.
        
        :param identifier: Relative path of the local file (e.g. 'Contracts/Acme_NDA.pdf') OR the Cloud File ID.
        :return: Structured markdown content of the parsed document.
        """
        if not self._is_authorized(__user__):
            return "Error: You are not authorized to read documents from the SharePoint repository."
            
        # Scenario A: Local Synced Folder
        if self.valves.SHAREPOINT_LOCAL_PATH:
            full_path = os.path.join(self.valves.SHAREPOINT_LOCAL_PATH, identifier)
            # If not found directly, check if identifier is just a filename
            if not os.path.exists(full_path):
                matches = glob.glob(os.path.join(self.valves.SHAREPOINT_LOCAL_PATH, "**", identifier), recursive=True)
                if matches:
                    full_path = matches[0]
                else:
                    return f"Error: Local file '{identifier}' not found under '{self.valves.SHAREPOINT_LOCAL_PATH}'."
            
            # Parse locally
            return self._parse_local_document(full_path)

        # Scenario B: Cloud SharePoint Connection
        elif self.valves.CLIENT_ID and self.valves.CLIENT_SECRET and self.valves.TENANT_ID:
            try:
                token = self._get_microsoft_token()
                headers = {"Authorization": f"Bearer {token}"}
                
                # Fetch drive item details to get name and download url
                url = f"https://graph.microsoft.com/v1.0/shares/{identifier}/driveItem" if identifier.startswith("u!") else f"https://graph.microsoft.com/v1.0/drive/items/{identifier}"
                response = requests.get(url, headers=headers)
                response.raise_for_status()
                item_info = response.json()
                
                filename = item_info.get("name")
                download_url = item_info.get("@microsoft.graph.downloadUrl")
                
                if not download_url:
                    return f"Error: Could not retrieve download URL for item ID '{identifier}'."
                
                # Download file to a local temp folder
                temp_dir = os.path.join(os.getcwd(), "temp_sharepoint_downloads")
                os.makedirs(temp_dir, exist_ok=True)
                local_path = os.path.join(temp_dir, filename)
                
                file_response = requests.get(download_url)
                file_response.raise_for_status()
                with open(local_path, "wb") as f:
                    f.write(file_response.content)
                
                # Parse the downloaded file locally
                parsed_content = self._parse_local_document(local_path)
                
                # Clean up local file after parsing
                if os.path.exists(local_path):
                    os.remove(local_path)
                    
                return f"--- Parsed SharePoint File: {filename} ---\n\n" + parsed_content
            except Exception as e:
                return f"Cloud SharePoint download/parse failed: {str(e)}"
        else:
            return "Error: SharePoint connection not configured."

    def _get_microsoft_token(self) -> str:
        import msal
        authority = f"https://login.microsoftonline.com/{self.valves.TENANT_ID}"
        app = msal.ConfidentialClientApplication(
            self.valves.CLIENT_ID,
            authority=authority,
            client_credential=self.valves.CLIENT_SECRET
        )
        result = app.acquire_token_for_client(scopes=["https://graph.microsoft.com/.default"])
        if "access_token" in result:
            return result["access_token"]
        else:
            raise Exception(f"Failed to acquire Microsoft token: {result.get('error_description')}")

    def _parse_local_document(self, filepath: str) -> str:
        ext = os.path.splitext(filepath)[1].lower()
        try:
            if ext == '.pdf':
                return self._parse_pdf(filepath)
            elif ext == '.docx':
                return self._parse_docx(filepath)
            elif ext == '.pptx':
                return self._parse_pptx(filepath)
            elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
                return self._parse_image(filepath)
            else:
                return f"Unsupported file type '{ext}' for local parsing."
        except Exception as e:
            return f"Error parsing file: {str(e)}"

    def _parse_pdf(self, filepath: str) -> str:
        import pdfplumber
        output = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                output.append(f"## Page {i + 1}\n")
                
                # Extract tables
                tables = page.extract_tables()
                if tables:
                    output.append("### Tables Extracted:\n")
                    for table in tables:
                        markdown_table = []
                        for r_idx, row in enumerate(table):
                            clean_row = [str(cell).strip().replace("\n", " ") if cell is not None else "" for cell in row]
                            markdown_table.append("| " + " | ".join(clean_row) + " |")
                            if r_idx == 0:
                                separator = "| " + " | ".join(["---"] * len(clean_row)) + " |"
                                markdown_table.append(separator)
                        output.append("\n".join(markdown_table) + "\n\n")
                
                # Extract text
                text = page.extract_text()
                if text:
                    output.append("### Text Content:\n")
                    output.append(text + "\n\n")
        return "\n".join(output)

    def _parse_docx(self, filepath: str) -> str:
        import docx
        doc = docx.Document(filepath)
        output = []
        for element in doc.element.body:
            if element.tag.endswith('p'):
                p = docx.text.paragraph.Paragraph(element, doc)
                if p.text.strip():
                    output.append(p.text + "\n")
            elif element.tag.endswith('tbl'):
                t = docx.table.Table(element, doc)
                markdown_table = []
                for r_idx, row in enumerate(t.rows):
                    row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    markdown_table.append("| " + " | ".join(row_cells) + " |")
                    if r_idx == 0:
                        separator = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                        markdown_table.append(separator)
                output.append("\n".join(markdown_table) + "\n\n")
        return "\n".join(output)

    def _parse_pptx(self, filepath: str) -> str:
        from pptx import Presentation
        prs = Presentation(filepath)
        output = []
        for i, slide in enumerate(prs.slides):
            output.append(f"## Slide {i + 1}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    output.append(shape.text.strip() + "\n")
                if shape.has_table:
                    markdown_table = []
                    for r_idx, row in enumerate(shape.table.rows):
                        row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        markdown_table.append("| " + " | ".join(row_cells) + " |")
                        if r_idx == 0:
                            separator = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                            markdown_table.append(separator)
                    output.append("\n".join(markdown_table) + "\n\n")
        return "\n".join(output)

    def _parse_image(self, filepath: str) -> str:
        from PIL import Image
        import pytesseract
        try:
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
            return f"### OCR Text Output:\n{text}"
        except Exception as e:
            return (
                f"Error performing OCR locally: {str(e)}.\n"
                "Please verify that 'tesseract-ocr' binary is installed on the computer hosting Open WebUI."
            )
