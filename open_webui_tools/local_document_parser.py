"""
title: Local Offline Document Parser
author: Advanced Agentic RAG Team
version: 1.0.0
description: Parses PDFs, Word docs, PowerPoint slides, and Scanned Images locally and completely offline. Preserves table structures, lists, and layout order as Markdown tables.
requirements: pdfplumber, python-docx, python-pptx, pytesseract, Pillow
"""

import os
from typing import Dict, List, Any

class Tools:
    def __init__(self):
        pass

    def parse_local_document(self, filename: str) -> str:
        """
        Parses a PDF, DOCX, PPTX, or image file locally and offline. Preserves document layouts, lists, and tabular data format as Markdown.
        
        :param filename: The name of the file to parse (e.g., 'report.pdf', 'table.xlsx').
        :return: Extracted document content as structured Markdown.
        """
        search_paths = [
            "/app/backend/data/uploads",
            "/mnt/uploads",
            "./backend/data/uploads",
            ".",
        ]
        
        filepath = None
        for path in search_paths:
            full_path = os.path.join(path, filename)
            if os.path.exists(full_path):
                filepath = full_path
                break
                
        if not filepath:
            # Check files in current dir
            if os.path.exists(filename):
                filepath = filename
            else:
                return f"Error: File '{filename}' not found."

        ext = os.path.splitext(filepath)[1].lower()
        
        try:
            if ext == '.pdf':
                return self._parse_pdf(filepath)
            elif ext == '.docx':
                return self._parse_docx(filepath)
            elif ext == '.pptx':
                return self._parse_pptx(filepath)
            elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
                return self._parse_image(filepath)
            else:
                return f"Unsupported file type '{ext}' for local parsing."
        except Exception as e:
            return f"Error parsing file: {str(e)}"

    def _parse_pdf(self, filepath: str) -> str:
        import pdfplumber
        
        output = []
        with pdfplumber.open(filepath) as pdf:
            for i, page in enumerate(pdf.pages):
                output.append(f"## Page {i + 1}\n")
                
                # 1. Extract tables
                tables = page.extract_tables()
                if tables:
                    output.append("### Tables Extracted:\n")
                    for table in tables:
                        markdown_table = []
                        for r_idx, row in enumerate(table):
                            clean_row = [str(cell).strip().replace("\n", " ") if cell is not None else "" for cell in row]
                            markdown_table.append("| " + " | ".join(clean_row) + " |")
                            if r_idx == 0:
                                separator = "| " + " | ".join(["---"] * len(clean_row)) + " |"
                                markdown_table.append(separator)
                        output.append("\n".join(markdown_table) + "\n\n")
                
                # 2. Extract text
                text = page.extract_text()
                if text:
                    output.append("### Text Content:\n")
                    output.append(text + "\n\n")
                    
        return "\n".join(output)

    def _parse_docx(self, filepath: str) -> str:
        import docx
        
        doc = docx.Document(filepath)
        output = []
        
        # Iterate sequentially through structural body elements
        for element in doc.element.body:
            if element.tag.endswith('p'):
                p = docx.text.paragraph.Paragraph(element, doc)
                if p.text.strip():
                    output.append(p.text + "\n")
            elif element.tag.endswith('tbl'):
                t = docx.table.Table(element, doc)
                markdown_table = []
                for r_idx, row in enumerate(t.rows):
                    row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                    markdown_table.append("| " + " | ".join(row_cells) + " |")
                    if r_idx == 0:
                        separator = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                        markdown_table.append(separator)
                output.append("\n".join(markdown_table) + "\n\n")
                
        return "\n".join(output)

    def _parse_pptx(self, filepath: str) -> str:
        from pptx import Presentation
        
        prs = Presentation(filepath)
        output = []
        
        for i, slide in enumerate(prs.slides):
            output.append(f"## Slide {i + 1}\n")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    output.append(shape.text.strip() + "\n")
                if shape.has_table:
                    markdown_table = []
                    for r_idx, row in enumerate(shape.table.rows):
                        row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        markdown_table.append("| " + " | ".join(row_cells) + " |")
                        if r_idx == 0:
                            separator = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                            markdown_table.append(separator)
                    output.append("\n".join(markdown_table) + "\n\n")
                    
        return "\n".join(output)

    def _parse_image(self, filepath: str) -> str:
        from PIL import Image
        import pytesseract
        
        try:
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
            return f"### OCR Text Output:\n{text}"
        except Exception as e:
            return (
                f"Error performing OCR locally: {str(e)}.\n"
                "Please verify that 'tesseract-ocr' binary is installed on the computer hosting Open WebUI "
                "(e.g., 'sudo apt-get install tesseract-ocr' on Linux)."
            )
