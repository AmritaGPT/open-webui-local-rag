"""
AmritaGPT Document Intelligence OpenAPI Server
Author: Advanced Agentic RAG Team
Version: 1.0.0
Description: A FastAPI server exposing layout-aware document search, retrieval, OCR, 
             and spreadsheet computations locally as a standard OpenAPI service.
"""

import os
import glob
import time
import base64
import io
import re
from typing import List, Dict, Any
from pydantic import BaseModel, Field
from fastapi import FastAPI, Query, HTTPException, Body
from fastapi.middleware.cors import CORSMiddleware
import uvicorn

# Data Processing Dependencies
import pandas as pd
import pdfplumber
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from PIL import Image
import pytesseract

# Initialize FastAPI App
app = FastAPI(
    title="AmritaGPT Document Intelligence API",
    description="Local layout-aware document search, parsing, and spreadsheet calculations.",
    version="1.0.0"
)

# Enable CORS for Open WebUI connectivity
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Default Synced Directory
DEFAULT_PATH = "/media/hirthikbalaji/AGPT DATA/SAMPLE"

def _resolve_safe_path(filename: str) -> str:
    """Sanitizes filename and resolves path, preventing directory traversal."""
    sanitized_filename = os.path.basename(filename)
    if "/" in filename or "\\" in filename:
        clean_parts = [os.path.basename(p) for p in re.split(r'[/\\]', filename) if p and p != '..']
        sanitized_filename = os.path.join(*clean_parts)
        
    full_path = os.path.normpath(os.path.join(DEFAULT_PATH, sanitized_filename))
    if not full_path.startswith(os.path.normpath(DEFAULT_PATH)):
        raise HTTPException(status_code=400, detail="Directory traversal attempt blocked.")
    return full_path

class SpreadsheetQueryRequest(BaseModel):
    filename: str = Field(..., description="The name of the spreadsheet file.")
    code: str = Field(..., description="The Python/Pandas code to execute. Store outcome in 'result' variable.")

@app.get("/")
def read_root():
    return {"message": "AmritaGPT Document Intelligence API is online and fully local."}

@app.get("/health")
def health_check():
    """Health check for local RAG document intelligence server."""
    dir_exists = os.path.exists(DEFAULT_PATH)
    return {
        "status": "healthy" if dir_exists else "degraded",
        "service": "AmritaGPT Document Intelligence API",
        "version": "1.1.0",
        "default_path": DEFAULT_PATH,
        "directory_accessible": dir_exists
    }


@app.get("/documents")
def list_all_documents():
    """
    Lists all documents available in the local repository recursively.
    """
    if not os.path.exists(DEFAULT_PATH):
        raise HTTPException(status_code=500, detail=f"Synced directory '{DEFAULT_PATH}' does not exist.")
        
    search_pattern = os.path.join(DEFAULT_PATH, "**", "*")
    files = glob.glob(search_pattern, recursive=True)
    
    results = []
    for filepath in files:
        if os.path.isfile(filepath):
            rel_path = os.path.relpath(filepath, DEFAULT_PATH)
            size_kb = os.path.getsize(filepath) / 1024
            mod_time = time.ctime(os.path.getmtime(filepath))
            results.append({
                "filename": os.path.basename(filepath),
                "path": rel_path,
                "size_kb": round(size_kb, 2),
                "last_modified": mod_time
            })
            
    return {"documents": results}

@app.get("/search")
def search_documents(query: str = Query(..., description="Keywords to search for"), search_contents: bool = True):
    """
    Searches the repository for documents matching keywords.
    Scans filenames and performs offline full-text content search.
    """
    if not os.path.exists(DEFAULT_PATH):
        raise HTTPException(status_code=500, detail="Repository path does not exist.")
        
    search_pattern = os.path.join(DEFAULT_PATH, "**", "*")
    all_items = glob.glob(search_pattern, recursive=True)
    
    matches = []
    # 1. Filename matching
    for filepath in all_items:
        if os.path.isfile(filepath):
            filename = os.path.basename(filepath)
            if query.lower() in filename.lower():
                matches.append({
                    "path": os.path.relpath(filepath, DEFAULT_PATH),
                    "reason": "Filename match",
                    "size_kb": round(os.path.getsize(filepath)/1024, 2)
                })
                
    # 2. Content search fallback (PDF, DOCX, TXT)
    if search_contents:
        for filepath in all_items:
            if os.path.isfile(filepath) and filepath not in [os.path.join(DEFAULT_PATH, m["path"]) for m in matches]:
                ext = os.path.splitext(filepath)[1].lower()
                try:
                    text_content = ""
                    if ext == '.txt':
                        with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                            text_content = f.read()
                    elif ext == '.pdf':
                        with pdfplumber.open(filepath) as pdf:
                            text_content = " ".join([page.extract_text() or "" for page in pdf.pages[:5]])
                    elif ext == '.docx':
                        doc = docx.Document(filepath)
                        text_content = " ".join([p.text for p in doc.paragraphs[:20]])
                        
                    if query.lower() in text_content.lower():
                        match_obj = re.search(rf"([^.?!]*?{re.escape(query)}[^.?!]*?[.?!])", text_content, re.IGNORECASE)
                        snippet = match_obj.group(1).strip() if match_obj else "Matches found in content."
                        matches.append({
                            "path": os.path.relpath(filepath, DEFAULT_PATH),
                            "reason": f"Content snippet: \"... {snippet} ...\"",
                            "size_kb": round(os.path.getsize(filepath)/1024, 2)
                        })
                except Exception:
                    continue
                    
    return {"results": matches[:10]}

@app.get("/parse")
def read_and_parse_document(filename: str = Query(..., description="The name or relative path of the file to parse")):
    """
    Parses a PDF, DOCX, PPTX, or Scanned Image locally.
    Maintains table structures, headings, and outlines sequentially in Markdown format.
    """
    try:
        filepath = _resolve_safe_path(filename)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
        
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="File not found.")
        
    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext in ['.txt', '.md', '.json', '.yaml', '.yml']:
            with open(filepath, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()
            return {"content": content}
        elif ext == '.csv':
            df = pd.read_csv(filepath)
            try:
                markdown_table = df.head(100).to_markdown(index=False)
            except Exception:
                # Fallback manual table parser if tabulate library is not present
                cols = list(df.columns)
                lines = [
                    "| " + " | ".join([str(c) for c in cols]) + " |",
                    "| " + " | ".join(["---"] * len(cols)) + " |"
                ]
                for _, row in df.head(100).iterrows():
                    lines.append("| " + " | ".join([str(val).replace("\n", " ") for val in row]) + " |")
                markdown_table = "\n".join(lines)
            return {"content": f"### CSV File Content (Top 100 Rows):\n\n{markdown_table}"}
        elif ext == '.pdf':
            output = []
            with pdfplumber.open(filepath) as pdf:
                for i, page in enumerate(pdf.pages):
                    output.append(f"## Page {i + 1}\n")
                    tables = page.extract_tables()
                    if tables:
                        output.append("### Tables Extracted:\n")
                        for table in tables:
                            markdown_table = []
                            for r_idx, row in enumerate(table):
                                clean_row = [str(cell).strip().replace("\n", " ") if cell is not None else "" for cell in row]
                                markdown_table.append("| " + " | ".join(clean_row) + " |")
                                if r_idx == 0:
                                    separator = "| " + " | ".join(["---"] * len(clean_row)) + " |"
                                    markdown_table.append(separator)
                            output.append("\n".join(markdown_table) + "\n\n")
                    text = page.extract_text()
                    if text:
                        output.append("### Text Content:\n" + text + "\n\n")
            return {"content": "\n".join(output)}
            
        elif ext == '.docx':
            doc = docx.Document(filepath)
            output = []
            for element in doc.element.body:
                if element.tag.endswith('p'):
                    p = docx.text.paragraph.Paragraph(element, doc)
                    if p.text.strip():
                        output.append(p.text + "\n")
                elif element.tag.endswith('tbl'):
                    t = docx.table.Table(element, doc)
                    markdown_table = []
                    for r_idx, row in enumerate(t.rows):
                        row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                        markdown_table.append("| " + " | ".join(row_cells) + " |")
                        if r_idx == 0:
                            separator = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                            markdown_table.append(separator)
                    output.append("\n".join(markdown_table) + "\n\n")
            return {"content": "\n".join(output)}
            
        elif ext == '.pptx':
            prs = Presentation(filepath)
            output = []
            for i, slide in enumerate(prs.slides):
                output.append(f"## Slide {i + 1}\n")
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        output.append(shape.text.strip() + "\n")
                    if shape.has_table:
                        markdown_table = []
                        for r_idx, row in enumerate(shape.table.rows):
                            row_cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
                            markdown_table.append("| " + " | ".join(row_cells) + " |")
                            if r_idx == 0:
                                separator = "| " + " | ".join(["---"] * len(row_cells)) + " |"
                                markdown_table.append(separator)
                        output.append("\n".join(markdown_table) + "\n\n")
            return {"content": "\n".join(output)}
            
        elif ext in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
            img = Image.open(filepath)
            text = pytesseract.image_to_string(img)
            return {"content": f"### OCR Text Output:\n{text}"}
        else:
            raise HTTPException(status_code=400, detail=f"Unsupported extension: {ext}")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error parsing file: {str(e)}")

@app.get("/spreadsheet/schema")
def get_spreadsheet_schema(filename: str = Query(..., description="Excel/CSV filename")):
    """
    Examines an Excel or CSV file structure, listing the sheet names, columns, and column data types.
    """
    try:
        filepath = _resolve_safe_path(filename)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
        
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="File not found.")
        
    try:
        if filepath.endswith('.csv'):
            df = pd.read_csv(filepath, nrows=5)
            cols = ", ".join([f"`{col}` ({df[col].dtype})" for col in df.columns])
            return {"type": "csv", "columns": cols}
        elif filepath.endswith(('.xlsx', '.xls')):
            xl = pd.ExcelFile(filepath)
            sheets = xl.sheet_names
            result = []
            for sheet in sheets:
                df = pd.read_excel(filepath, sheet_name=sheet, nrows=5)
                cols = ", ".join([f"`{col}` ({df[col].dtype})" for col in df.columns])
                result.append({
                    "sheet_name": sheet,
                    "columns": cols
                })
            return {"type": "excel", "sheets": result}
        else:
            raise HTTPException(status_code=400, detail="Unsupported spreadsheet type.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error reading schema: {str(e)}")

@app.post("/spreadsheet/query")
def query_spreadsheet(payload: SpreadsheetQueryRequest):
    """
    Runs Python Pandas code on a local spreadsheet to execute calculations, filters, or category groups.
    """
    try:
        filepath = _resolve_safe_path(payload.filename)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
        
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="File not found.")
        
    try:
        if filepath.endswith('.csv'):
            df = pd.read_csv(filepath)
            sheets = {"default": df}
        elif filepath.endswith(('.xlsx', '.xls')):
            sheets = pd.read_excel(filepath, sheet_name=None)
            df = list(sheets.values())[0] if len(sheets) == 1 else None
        else:
            raise HTTPException(status_code=400, detail="Unsupported file format.")
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error loading file: {str(e)}")
        
    local_vars = {
        'sheets': sheets,
        'df': df,
        'pd': pd,
        'result': None
    }
    
    try:
        exec(payload.code, {}, local_vars)
        result = local_vars.get('result')
        if result is None:
            raise HTTPException(status_code=400, detail="Code ran successfully, but 'result' variable was not assigned.")
        return {"result": str(result)}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error executing calculation: {str(e)}")

@app.get("/pdf/render")
def render_pdf_page_as_image(filename: str = Query(..., description="PDF filename"), page_num: int = 1):
    """
    Converts a single PDF page into a base64 PNG data URL.
    """
    try:
        filepath = _resolve_safe_path(filename)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
        
    if not os.path.exists(filepath):
        raise HTTPException(status_code=404, detail="File not found.")
        
    try:
        doc = fitz.open(filepath)
        total_pages = len(doc)
        idx = max(1, min(total_pages, page_num)) - 1
        
        page = doc.load_page(idx)
        zoom = 300 / 72
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        
        img_data = pix.tobytes("png")
        img = Image.open(io.BytesIO(img_data))
        
        buffered = io.BytesIO()
        img.save(buffered, format="PNG")
        img_str = base64.b64encode(buffered.getvalue()).decode("utf-8")
        
        return {
            "page_num": idx + 1,
            "image_data_url": f"data:image/png;base64,{img_str}"
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error rendering page: {str(e)}")

if __name__ == "__main__":
    uvicorn.run(app, host="0.0.0.0", port=8000)
